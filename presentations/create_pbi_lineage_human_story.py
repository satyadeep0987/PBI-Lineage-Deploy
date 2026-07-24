"""Create the two-slide PBI Lineage Explorer executive presentation."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "presentations" / "PBI_Lineage_Explorer_User_Value_Architecture.pptx"
PREVIEW_DIR = ROOT / "presentations" / "previews"

SLIDE_W = 13.333
SLIDE_H = 7.5
PREVIEW_SCALE = 120

NAVY = "10213D"
INK = "17233A"
MUTED = "5E6D85"
BLUE = "2563EB"
BLUE_LIGHT = "EAF2FF"
GREEN = "059669"
GREEN_LIGHT = "E8F7F1"
ORANGE = "EA580C"
ORANGE_LIGHT = "FFF0E8"
TEAL = "0F766E"
TEAL_LIGHT = "E8F7F5"
PAPER = "FFFFFF"
CANVAS = "F5F7FB"
LINE = "D7DFEA"
SOFT = "EEF2F7"

FONT_REGULAR = Path(r"C:\Windows\Fonts\segoeui.ttf")
FONT_BOLD = Path(r"C:\Windows\Fonts\segoeuib.ttf")


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor.from_string(value)


def pil_color(hex_color: str) -> str:
    return f"#{hex_color.lstrip('#')}"


class SlideCanvas:
    """Draw matching editable PowerPoint shapes and a PNG preview."""

    def __init__(self, prs: Presentation):
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.image = Image.new(
            "RGB",
            (round(SLIDE_W * PREVIEW_SCALE), round(SLIDE_H * PREVIEW_SCALE)),
            pil_color(CANVAS),
        )
        self.draw = ImageDraw.Draw(self.image)
        background = self.slide.background.fill
        background.solid()
        background.fore_color.rgb = rgb(CANVAS)

    @staticmethod
    def _px(value: float) -> int:
        return round(value * PREVIEW_SCALE)

    @staticmethod
    def _font(size: float, bold: bool = False) -> ImageFont.FreeTypeFont:
        path = FONT_BOLD if bold else FONT_REGULAR
        return ImageFont.truetype(str(path), round(size * PREVIEW_SCALE / 72))

    def rect(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: str,
        line: str | None = None,
        line_width: float = 1,
        rounded: bool = True,
    ):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = self.slide.shapes.add_shape(
            shape_type,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line:
            shape.line.color.rgb = rgb(line)
            shape.line.width = Pt(line_width)
        else:
            shape.line.fill.background()

        bounds = [
            self._px(x),
            self._px(y),
            self._px(x + w),
            self._px(y + h),
        ]
        if rounded:
            self.draw.rounded_rectangle(
                bounds,
                radius=max(2, self._px(0.08)),
                fill=pil_color(fill),
                outline=pil_color(line) if line else None,
                width=max(1, round(line_width * PREVIEW_SCALE / 72)),
            )
        else:
            self.draw.rectangle(
                bounds,
                fill=pil_color(fill),
                outline=pil_color(line) if line else None,
                width=max(1, round(line_width * PREVIEW_SCALE / 72)),
            )
        return shape

    def line(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        color: str = LINE,
        width: float = 1,
    ):
        connector = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        connector.line.color.rgb = rgb(color)
        connector.line.width = Pt(width)
        self.draw.line(
            [
                self._px(x1),
                self._px(y1),
                self._px(x2),
                self._px(y2),
            ],
            fill=pil_color(color),
            width=max(1, round(width * PREVIEW_SCALE / 72)),
        )
        return connector

    def arrow(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        color: str = BLUE,
        direction: str = "right",
    ):
        shape_type = {
            "right": MSO_SHAPE.RIGHT_ARROW,
            "left": MSO_SHAPE.LEFT_ARROW,
            "down": MSO_SHAPE.DOWN_ARROW,
        }[direction]
        shape = self.slide.shapes.add_shape(
            shape_type,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()

        left = self._px(x)
        top = self._px(y)
        right = self._px(x + w)
        bottom = self._px(y + h)
        middle = (top + bottom) // 2
        notch = round((right - left) * 0.34)
        inset = round((bottom - top) * 0.25)
        if direction == "right":
            points = [
                (left, top + inset),
                (right - notch, top + inset),
                (right - notch, top),
                (right, middle),
                (right - notch, bottom),
                (right - notch, bottom - inset),
                (left, bottom - inset),
            ]
        elif direction == "left":
            points = [
                (right, top + inset),
                (left + notch, top + inset),
                (left + notch, top),
                (left, middle),
                (left + notch, bottom),
                (left + notch, bottom - inset),
                (right, bottom - inset),
            ]
        else:
            vertical_notch = round((bottom - top) * 0.34)
            horizontal_inset = round((right - left) * 0.25)
            center = (left + right) // 2
            points = [
                (left + horizontal_inset, top),
                (right - horizontal_inset, top),
                (right - horizontal_inset, bottom - vertical_notch),
                (right, bottom - vertical_notch),
                (center, bottom),
                (left, bottom - vertical_notch),
                (left + horizontal_inset, bottom - vertical_notch),
            ]
        self.draw.polygon(points, fill=pil_color(color))
        return shape

    def text(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        size: float,
        color: str = INK,
        bold: bool = False,
        align: str = "left",
        valign: str = "top",
        margin: float = 0.04,
    ):
        box = self.slide.shapes.add_textbox(
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }[valign]
        paragraph = frame.paragraphs[0]
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[align]
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = value
        run.font.name = "Segoe UI"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)

        font = self._font(size, bold)
        max_width = max(1, self._px(w - (margin * 2)))
        wrapped_lines: list[str] = []
        for explicit_line in value.splitlines() or [""]:
            words = explicit_line.split()
            if not words:
                wrapped_lines.append("")
                continue
            current = words[0]
            for word in words[1:]:
                candidate = f"{current} {word}"
                if self.draw.textlength(candidate, font=font) <= max_width:
                    current = candidate
                else:
                    wrapped_lines.append(current)
                    current = word
            wrapped_lines.append(current)
        wrapped = "\n".join(wrapped_lines)
        spacing = max(2, round(size * 0.28 * PREVIEW_SCALE / 72))
        bbox = self.draw.multiline_textbbox(
            (0, 0),
            wrapped,
            font=font,
            spacing=spacing,
            align=align,
        )
        text_w = bbox[2] - bbox[0]
        text_h = bbox[3] - bbox[1]
        left = self._px(x + margin)
        top = self._px(y + margin)
        content_w = self._px(w - (margin * 2))
        content_h = self._px(h - (margin * 2))
        if align == "center":
            left += max(0, (content_w - text_w) // 2)
        elif align == "right":
            left += max(0, content_w - text_w)
        if valign == "middle":
            top += max(0, (content_h - text_h) // 2)
        elif valign == "bottom":
            top += max(0, content_h - text_h)
        self.draw.multiline_text(
            (left, top),
            wrapped,
            font=font,
            fill=pil_color(color),
            spacing=spacing,
            align=align,
        )
        return box

    def label_box(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        subtitle: str,
        fill: str,
        accent: str,
        title_size: float = 12,
    ):
        self.rect(x, y, w, h, fill, LINE, 0.8)
        self.rect(x, y, 0.07, h, accent, None, rounded=False)
        self.text(x + 0.2, y + 0.12, w - 0.32, 0.26, title, title_size, INK, True)
        self.text(x + 0.2, y + 0.44, w - 0.32, h - 0.53, subtitle, 8.5, MUTED)

    def save_preview(self, name: str):
        PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
        path = PREVIEW_DIR / name
        self.image.save(path, "PNG")
        return path


def add_header(canvas: SlideCanvas, section: str, slide_number: str):
    canvas.rect(0.55, 0.34, 0.44, 0.40, BLUE, None)
    canvas.text(0.55, 0.34, 0.44, 0.40, "PBI", 9, PAPER, True, "center", "middle", 0)
    canvas.text(1.08, 0.33, 2.4, 0.42, "Lineage Explorer", 15, NAVY, True, valign="middle")
    canvas.text(9.65, 0.37, 2.85, 0.30, section, 8.5, BLUE, True, "right", "middle")
    canvas.text(12.53, 0.37, 0.25, 0.30, slide_number, 8.5, MUTED, True, "right", "middle")
    canvas.line(0.55, 0.88, 12.78, 0.88, LINE, 0.8)


def add_footer(canvas: SlideCanvas, slide_number: str):
    canvas.text(
        0.58,
        7.14,
        5.2,
        0.22,
        "PBI LINEAGE EXPLORER  |  POWER BI + FABRIC + XMLA + SNOWFLAKE",
        7.5,
        MUTED,
        True,
        valign="middle",
    )
    canvas.text(12.38, 7.14, 0.38, 0.22, slide_number, 7.5, MUTED, True, "right", "middle")


def add_path_node(
    canvas: SlideCanvas,
    x: float,
    y: float,
    title: str,
    code: str,
    fill: str,
    accent: str,
):
    canvas.rect(x, y, 1.28, 0.92, fill, LINE, 0.8)
    canvas.rect(x + 0.12, y + 0.14, 0.32, 0.32, accent, None)
    canvas.text(x + 0.12, y + 0.14, 0.32, 0.32, code, 8.5, PAPER, True, "center", "middle", 0)
    canvas.text(x + 0.12, y + 0.55, 1.04, 0.24, title, 9.2, INK, True, "center", "middle", 0)


def build_slide_one(prs: Presentation) -> SlideCanvas:
    canvas = SlideCanvas(prs)
    add_header(canvas, "01  /  USER VALUE", "01")
    canvas.text(
        0.58,
        1.12,
        12.0,
        0.50,
        "Get lineage answers in minutes, not hours",
        27,
        NAVY,
        True,
        valign="middle",
    )
    canvas.text(
        0.58,
        1.68,
        11.8,
        0.38,
        "Instead of opening reports one by one, a BI user searches once and sees the connected report, model, measure, visual and source details.",
        13,
        MUTED,
        valign="middle",
    )

    canvas.rect(0.58, 2.28, 3.12, 3.92, NAVY, None)
    canvas.text(0.84, 2.53, 2.60, 0.25, "THE OLD WAY", 8.5, "8FB7FF", True)
    canvas.text(
        0.84,
        2.91,
        2.52,
        1.58,
        "\"Open every report, inspect every model, and still wonder what changed.\"",
        17,
        PAPER,
        True,
        valign="middle",
    )
    canvas.line(0.84, 4.56, 3.38, 4.56, "40506A", 0.8)
    canvas.text(
        0.84,
        4.72,
        2.50,
        0.38,
        "The user has to jump between:",
        9.5,
        "D9E3F1",
    )
    for x, y, label in [
        (0.84, 5.20, "POWER BI"),
        (2.05, 5.20, "FABRIC"),
        (0.84, 5.60, "XMLA"),
        (2.05, 5.60, "SNOWFLAKE"),
    ]:
        canvas.rect(x, y, 1.08, 0.29, "243653", "40506A", 0.6)
        canvas.text(x, y, 1.08, 0.29, label, 7.3, PAPER, True, "center", "middle", 0)

    canvas.text(4.02, 2.32, 4.70, 0.34, "The app does the report-by-report work", 14, NAVY, True)
    canvas.text(
        4.02,
        2.69,
        4.65,
        0.42,
        "It inventories authorized assets, joins metadata by stable IDs and traces impact across BI and source layers.",
        9.5,
        MUTED,
    )

    add_path_node(canvas, 4.04, 3.27, "Report", "R", BLUE_LIGHT, BLUE)
    add_path_node(canvas, 5.67, 3.27, "Visual", "V", GREEN_LIGHT, GREEN)
    add_path_node(canvas, 7.30, 3.27, "Measure", "M", ORANGE_LIGHT, ORANGE)
    canvas.arrow(5.37, 3.61, 0.22, 0.18, BLUE)
    canvas.arrow(7.00, 3.61, 0.22, 0.18, BLUE)

    add_path_node(canvas, 4.04, 4.68, "Raw columns", "01", GREEN_LIGHT, GREEN)
    add_path_node(canvas, 5.67, 4.68, "Snowflake", "DB", TEAL_LIGHT, TEAL)
    add_path_node(canvas, 7.30, 4.68, "Model", "S", BLUE_LIGHT, BLUE)
    canvas.arrow(7.80, 4.28, 0.28, 0.30, BLUE, "down")
    canvas.arrow(7.00, 5.02, 0.22, 0.18, TEAL, "left")
    canvas.arrow(5.37, 5.02, 0.22, 0.18, TEAL, "left")

    canvas.rect(4.02, 5.82, 4.57, 0.38, SOFT, None)
    canvas.text(
        4.02,
        5.82,
        4.57,
        0.38,
        "ONE SEARCH  |  MINUTES TO AN EVIDENCE-BACKED ANSWER",
        8,
        NAVY,
        True,
        "center",
        "middle",
        0,
    )

    canvas.text(9.05, 2.32, 3.60, 0.34, "What people gain", 14, NAVY, True)
    canvas.label_box(
        9.05,
        2.88,
        3.60,
        0.84,
        "Minutes, not hours",
        "Ask once instead of checking every report and semantic model manually.",
        PAPER,
        BLUE,
    )
    canvas.label_box(
        9.05,
        3.89,
        3.60,
        0.84,
        "Complete impact view",
        "See affected reports, measures, visuals, tables and Snowflake sources.",
        PAPER,
        GREEN,
    )
    canvas.label_box(
        9.05,
        4.90,
        3.60,
        0.84,
        "Clean output",
        "Download the actual PowerAI response as simple Markdown.",
        PAPER,
        ORANGE,
    )
    canvas.text(
        9.08,
        5.89,
        3.55,
        0.30,
        "The app explains evidence; it does not invent lineage.",
        8.7,
        MUTED,
        True,
        "center",
        "middle",
    )

    canvas.rect(0.58, 6.55, 12.07, 0.46, NAVY, None)
    canvas.text(
        0.78,
        6.55,
        11.67,
        0.46,
        "OUTCOME: FEWER MANUAL CHECKS. FASTER CONFIDENCE. SAFER BI CHANGES.",
        10,
        PAPER,
        True,
        "center",
        "middle",
        0,
    )
    add_footer(canvas, "01")
    canvas.slide.notes_slide.notes_text_frame.text = (
        "Open with the old manual workflow: report owners spend hours opening reports, "
        "checking semantic models and tracing source objects. The app compresses that into "
        "one search and one evidence trail. Emphasize minutes instead of report-by-report "
        "inspection, then close with the user outcomes: faster confidence, safer changes "
        "and a clean PowerAI answer that can be shared."
    )
    return canvas


def add_architecture_system(
    canvas: SlideCanvas,
    y: float,
    code: str,
    title: str,
    subtitle: str,
    fill: str,
    accent: str,
):
    canvas.rect(6.86, y, 3.12, 0.93, fill, LINE, 0.8)
    canvas.rect(7.05, y + 0.18, 0.48, 0.48, accent, None)
    canvas.text(7.05, y + 0.18, 0.48, 0.48, code, 8, PAPER, True, "center", "middle", 0)
    canvas.text(7.68, y + 0.14, 2.08, 0.26, title, 10.2, INK, True)
    canvas.text(7.68, y + 0.45, 2.08, 0.30, subtitle, 7.8, MUTED)


def add_change_chip(
    canvas: SlideCanvas,
    x: float,
    title: str,
    subtitle: str,
    accent: str,
):
    canvas.rect(x, 6.14, 2.33, 0.72, PAPER, LINE, 0.7)
    canvas.rect(x, 6.14, 0.06, 0.72, accent, None, rounded=False)
    canvas.text(x + 0.16, 6.25, 2.02, 0.18, title, 7.6, accent, True)
    canvas.text(x + 0.16, 6.47, 2.02, 0.24, subtitle, 8.2, INK, True)


def build_slide_two(prs: Presentation) -> SlideCanvas:
    canvas = SlideCanvas(prs)
    add_header(canvas, "02  /  ARCHITECTURE + BI USER FLOW", "02")
    canvas.text(
        0.58,
        1.12,
        12.0,
        0.50,
        "How each layer helps the BI user interact with lineage",
        27,
        NAVY,
        True,
        valign="middle",
    )
    canvas.text(
        0.58,
        1.68,
        11.9,
        0.38,
        "The BI user stays in one Streamlit app while app-owned connectors collect evidence from Power BI, Fabric, XMLA, Snowflake and PowerAI.",
        12.5,
        MUTED,
        valign="middle",
    )

    canvas.text(0.58, 2.31, 1.32, 0.23, "PEOPLE", 8, BLUE, True, "center")
    canvas.rect(0.58, 2.67, 1.32, 1.48, BLUE_LIGHT, "BCD1F5", 0.8)
    canvas.rect(0.98, 2.93, 0.52, 0.52, BLUE, None)
    canvas.text(0.98, 2.93, 0.52, 0.52, "U", 12, PAPER, True, "center", "middle", 0)
    canvas.text(0.64, 3.56, 1.20, 0.23, "BI user / owner", 8.1, INK, True, "center")
    canvas.text(0.68, 3.84, 1.12, 0.20, "Search, explore, ask", 7.3, MUTED, False, "center")
    canvas.arrow(2.04, 3.24, 0.30, 0.28, BLUE)

    canvas.text(2.50, 2.31, 3.70, 0.23, "PBI LINEAGE EXPLORER", 8, BLUE, True, "center")
    canvas.rect(2.50, 2.67, 3.70, 2.98, PAPER, BLUE, 1.2)
    canvas.rect(2.50, 2.67, 3.70, 0.44, NAVY, None, rounded=False)
    canvas.text(
        2.66,
        2.67,
        3.38,
        0.44,
        "STREAMLIT + PYTHON ORCHESTRATION",
        9,
        PAPER,
        True,
        "center",
        "middle",
        0,
    )
    for y, code, title, subtitle, accent in [
        (3.26, "01", "Search and navigation", "Home, Explore, Report Lineage, Table/Measure Impact", BLUE),
        (3.88, "02", "Session and permissions", "Signed-in identity scopes every request", GREEN),
        (4.50, "03", "Lineage engine", "Joins report, visual, DAX, model and source metadata", ORANGE),
        (5.12, "04", "PowerAI and export", "Explains evidence; Markdown download is answer-only", TEAL),
    ]:
        canvas.rect(2.70, y, 0.38, 0.30, accent, None)
        canvas.text(2.70, y, 0.38, 0.30, code, 7.4, PAPER, True, "center", "middle", 0)
        canvas.text(3.19, y - 0.02, 2.75, 0.24, title, 9.2, INK, True)
        canvas.text(3.19, y + 0.24, 2.75, 0.26, subtitle, 7.4, MUTED)
        if y < 5.12:
            canvas.line(2.70, y + 0.52, 6.00, y + 0.52, LINE, 0.6)

    canvas.arrow(6.35, 3.75, 0.30, 0.28, TEAL)
    canvas.text(6.86, 2.31, 3.12, 0.23, "PLATFORM METADATA", 8, TEAL, True, "center")
    add_architecture_system(
        canvas,
        2.67,
        "API",
        "Power BI REST + Fabric",
        "Workspaces, reports, pages, visuals and definitions",
        BLUE_LIGHT,
        BLUE,
    )
    add_architecture_system(
        canvas,
        3.77,
        "X",
        "Windows XMLA path",
        "Tables, columns, measures, DAX and dependencies",
        ORANGE_LIGHT,
        ORANGE,
    )
    add_architecture_system(
        canvas,
        4.87,
        "SF",
        "Snowflake lineage",
        "Tables, columns, transformations and raw sources",
        TEAL_LIGHT,
        TEAL,
    )

    canvas.arrow(10.13, 3.75, 0.30, 0.28, NAVY)
    canvas.text(10.55, 2.31, 2.18, 0.23, "BI USER SEES", 8, NAVY, True, "center")
    canvas.rect(10.55, 2.67, 2.18, 2.02, NAVY, None)
    canvas.text(10.78, 2.94, 1.72, 0.24, "ASK ONCE", 8.2, "8FB7FF", True, "center")
    canvas.text(10.78, 3.35, 1.72, 0.28, "Reports affected", 11, PAPER, True, "center")
    canvas.text(10.78, 3.72, 1.72, 0.28, "Model objects used", 11, PAPER, True, "center")
    canvas.text(10.78, 4.09, 1.72, 0.28, "Source tables feeding it", 10.5, PAPER, True, "center")

    canvas.rect(10.55, 4.88, 2.18, 0.77, ORANGE_LIGHT, ORANGE, 1.0)
    canvas.text(10.72, 4.98, 1.84, 0.18, "POWERAI ASSISTANT", 7.8, ORANGE, True, "center")
    canvas.text(
        10.72,
        5.20,
        1.84,
        0.30,
        "Floating modal answers using read-only app evidence",
        8.2,
        INK,
        True,
        "center",
    )

    add_change_chip(canvas, 0.58, "AUTHENTICATED USER", "Only authorized assets", BLUE)
    add_change_chip(canvas, 3.04, "CONNECTED METADATA", "REST + Fabric + XMLA + Snowflake", ORANGE)
    add_change_chip(canvas, 5.50, "POWERAI MODAL", "Ask without leaving page", TEAL)
    add_change_chip(canvas, 7.96, "ANSWER EXPORT", "Clean .md output only", GREEN)
    canvas.text(
        10.55,
        6.15,
        2.18,
        0.70,
        "Every answer is tied to app-collected metadata;\nPowerAI explains but does not change data.",
        8,
        MUTED,
        True,
        "center",
        "middle",
    )

    add_footer(canvas, "02")
    canvas.slide.notes_slide.notes_text_frame.text = (
        "Walk from left to right. The BI user stays in one Streamlit experience and interacts "
        "through Home, Explore, Report Lineage, impact pages and the floating PowerAI modal. "
        "The middle layer keeps the authenticated session, joins metadata by stable IDs and "
        "runs lineage/impact analysis. The platform boxes explain each role: REST and Fabric "
        "provide inventory and report definitions, XMLA provides semantic-model detail, "
        "Snowflake provides physical lineage, and PowerAI explains the app-owned evidence "
        "without performing write actions."
    )
    return canvas


def validate_deck(prs: Presentation):
    if len(prs.slides) != 2:
        raise RuntimeError(f"Expected 2 slides, found {len(prs.slides)}")
    for slide_number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise RuntimeError(f"Slide {slide_number} has a shape outside the canvas")
            if shape.left + shape.width > prs.slide_width:
                raise RuntimeError(f"Slide {slide_number} has a shape beyond the right edge")
            if shape.top + shape.height > prs.slide_height:
                raise RuntimeError(f"Slide {slide_number} has a shape beyond the bottom edge")


def main():
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "PBI Lineage Explorer: User Value and Architecture"
    prs.core_properties.subject = "Two-slide user value and architecture overview"
    prs.core_properties.author = "PBI Lineage Explorer"
    prs.core_properties.keywords = "Power BI, Snowflake, lineage, architecture, impact analysis"

    slide_one = build_slide_one(prs)
    slide_two = build_slide_two(prs)
    validate_deck(prs)
    prs.save(OUTPUT_PATH)
    preview_one = slide_one.save_preview("slide_01_user_value.png")
    preview_two = slide_two.save_preview("slide_02_architecture_flow.png")

    print(f"Created: {OUTPUT_PATH}")
    print(f"Preview: {preview_one}")
    print(f"Preview: {preview_two}")


if __name__ == "__main__":
    main()
